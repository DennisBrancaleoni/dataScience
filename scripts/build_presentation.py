"""Build F1 presentation matching style of F1_New_Presentation.pptx"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches as In
import os

BASE = os.path.dirname(os.path.abspath(__file__))

# ── colours ───────────────────────────────────────────────────────────────────
RED   = RGBColor(0xCC, 0x00, 0x00)
DARK  = RGBColor(0x1A, 0x1A, 0x1A)
CARD  = RGBColor(0x21, 0x21, 0x21)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
MUTED = RGBColor(0xAA, 0xAA, 0xAA)

# ── helpers ───────────────────────────────────────────────────────────────────
def I(inches): return Inches(inches)

def set_bg(slide, colour):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = colour

def add_rect(slide, l, t, w, h, fill_rgb):
    shape = slide.shapes.add_shape(1, I(l), I(t), I(w), I(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    shape.line.fill.background()
    return shape

def add_text(slide, text, l, t, w, h, size, bold=False, colour=WHITE, align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(I(l), I(t), I(w), I(h))
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = colour
    return txb

def red_bar(slide, title_text, title_size=22):
    """Top red bar + white title"""
    add_rect(slide, 0, 0, 10, 1.0, RED)
    add_text(slide, title_text, 0.4, 0.1, 9.2, 0.8,
             title_size, bold=True, colour=WHITE, align=PP_ALIGN.LEFT)

def info_cards_right(slide, cards):
    """4 stacked info cards on the right (slide 3 layout)"""
    card_h, gap = 0.984, 0.0
    top_start = 1.1
    for i, (hdr, sub) in enumerate(cards):
        t = top_start + i * (card_h + gap)
        add_rect(slide, 6.048, t, 3.652, card_h, CARD)
        add_text(slide, hdr, 6.15, t + 0.05, 3.45, 0.42,
                 9.5, bold=True, colour=WHITE, wrap=True)
        add_text(slide, sub, 6.15, t + 0.48, 3.45, 0.45,
                 8, bold=False, colour=MUTED, wrap=True)

def info_cards_bottom(slide, cards):
    """4 horizontal info cards at the bottom (slide 4 layout)"""
    card_w, gap = 2.297, 0.1
    left_start = 0.3
    t = 4.72
    for i, (hdr, sub) in enumerate(cards):
        l = left_start + i * (card_w + gap)
        add_rect(slide, l, t, card_w, 0.78, CARD)
        add_text(slide, hdr, l + 0.1, t + 0.04, card_w - 0.2, 0.36,
                 8.5, bold=True, colour=WHITE, wrap=True)
        add_text(slide, sub, l + 0.1, t + 0.42, card_w - 0.2, 0.3,
                 7.5, bold=False, colour=MUTED, wrap=True)

def add_image_left(slide, img_path, w=5.5, h=4.25):
    """Image on the left (slide 3 layout)"""
    slide.shapes.add_picture(img_path, I(0.3), I(1.1), I(w), I(h))

def add_image_wide(slide, img_path, w=9.4, h=3.5):
    """Wide image (slide 4 layout) — centred"""
    l = (10 - w) / 2
    slide.shapes.add_picture(img_path, I(l), I(1.05), I(w), I(h))


# ── build presentation ────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = I(10)
prs.slide_height = I(5.625)
blank = prs.slide_layouts[6]   # completely blank layout


# ════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
set_bg(sl, DARK)

# Left red stripe
add_rect(sl, 0, 0, 0.18, 5.625, RED)

# "F1" big
add_text(sl, "F1", 0.35, 0.35, 2.2, 1.7, 72, bold=True, colour=RED)

# Main title
add_text(sl, "Formula 1 World Championship Analysis",
         0.35, 1.85, 7.2, 0.7, 26, bold=True, colour=WHITE)

# Subtitle
add_text(sl, "Dataset (1950 – 2024)",
         0.35, 2.65, 5.0, 0.6, 26, bold=False, colour=RED)

# Divider
add_rect(sl, 0.35, 3.42, 4.2, 0.045, RGBColor(0xC0, 0xC0, 0xC0))

# Presented by
add_text(sl, "Presentation by:", 0.35, 3.52, 4.0, 0.35,
         10, colour=MUTED)
add_text(sl, "Dennis Brancaleoni  ·  Davide Bras Ferrari",
         0.35, 3.88, 6.0, 0.4, 13, bold=True, colour=WHITE)

# Course label
add_text(sl, "Data Science Course", 0.35, 4.4, 4.0, 0.35,
         9, colour=MUTED)


# ════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Dataset Overview (row layout)
# ════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
set_bg(sl, DARK)
red_bar(sl, "The Dataset: Formula 1 World Championship (1950–2024)")

rows = [
    ("01", "75 Years of Racing History",
           "From the inaugural 1950 British Grand Prix to the 2024 Abu Dhabi GP — a complete record of every race"),
    ("02", "Scale: 26,759 Race Entries — 802 Drivers — 211 Constructors",
           "Linked across 7 CSV tables: results, races, drivers, constructors, circuits, qualifying, pit_stops"),
    ("03", "Source: Ergast Motor Racing Developer API (extended to 2024)",
           "Flat joined file — mainDataset.csv — used for all analyses; no manual manipulation"),
    ("04", "Why It Matters",
           "Four distinct technical regulation eras (V10 → V8 → V6 Hybrid → Ground Effect) create natural experiments in the data"),
]
row_h = 0.85
gap   = 0.09
top0  = 1.1
for i, (num, hdr, sub) in enumerate(rows):
    t = top0 + i * (row_h + gap)
    add_rect(sl, 0.4, t, 9.2, row_h, CARD)
    add_rect(sl, 0.4, t, 0.7, row_h, RED)   # accent block
    add_text(sl, num,  0.4,  t,        0.7, row_h, 14, bold=True, colour=WHITE, align=PP_ALIGN.CENTER)
    add_text(sl, hdr, 1.25, t + 0.06,  8.2, 0.34, 11, bold=True,  colour=WHITE)
    add_text(sl, sub, 1.25, t + 0.43,  8.2, 0.35,  9, bold=False, colour=MUTED)


# ════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Key Statistics & Calendar Growth
# ════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
set_bg(sl, DARK)
red_bar(sl, "Key Statistics: 75 Seasons at a Glance")

# KPI cards image (top half, wide)
sl.shapes.add_picture(os.path.join(BASE, "pres_kpi.png"),
                      I(0.3), I(1.05), I(9.4), I(1.85))

# Races-per-season chart (bottom half)
sl.shapes.add_picture(os.path.join(BASE, "pres_races_season.png"),
                      I(0.3), I(3.0), I(9.4), I(2.55))


# ════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Geographic Distribution
# ════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
set_bg(sl, DARK)
red_bar(sl, "Analysis 1: Geographic Distribution of F1 Races (1950–2024)")

add_image_wide(sl, os.path.join(BASE, "pres_geographic.png"), w=9.4, h=3.55)

info_cards_bottom(sl, [
    ("UK leads all-time with 75 races",
     "Silverstone has hosted every season since 1950"),
    ("Italy & Germany follow",
     "Ferrari's home and Europe's engineering heartland"),
    ("Expansion from 1990s onward",
     "Asia, Middle East and Americas entered the calendar"),
    ("2010s–2024: Middle East growth",
     "Bahrain, Abu Dhabi and Saudi Arabia became regular fixtures"),
])


# ════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Race Outcome Distribution
# ════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
set_bg(sl, DARK)
red_bar(sl, "Analysis 2: Race Outcome Distribution — All 26,759 Entries")

add_image_left(sl, os.path.join(BASE, "pres_donut.png"), w=5.6, h=4.4)

info_cards_right(sl, [
    ("~50% of entries: Classified Finish",
     "Half of all starters historically crossed the finish line"),
    ("Mechanical failures: ~20%",
     "Early-era unreliability; proportion halved after 2000"),
    ("Accidents / Collisions: ~10%",
     "Small by count but large in race-impact (safety cars, DNFs)"),
    ("DNS / DNQ: pre-1996 era",
     "Grids >26 cars required pre-qualifying in the 1980s–1990s"),
])


# ════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Pit Stop Count by Start & Finish Position
# ════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
set_bg(sl, DARK)
red_bar(sl, "Analysis 3: Pit Stop Strategy — Start vs. Finish Position")

add_image_left(sl, os.path.join(BASE, "pres_pitstop_heatmap.png"), w=5.6, h=4.35)

info_cards_right(sl, [
    ("Front runners use fewer stops (~1.5)",
     "Leaders control strategy and avoid traffic during pit windows"),
    ("Mid-field and back: up to 2.5+ stops",
     "Aggressive strategies used to gain positions from P10–P20"),
    ("No stops alone cure a poor grid position",
     "Outright pace dominates: extra stops rarely flip P15 → P5"),
    ("Data covers post-2010 refuelling-ban era",
     "Stop counts became strategically meaningful only after 2010"),
])


# ════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Driver Win Rate
# ════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
set_bg(sl, DARK)
red_bar(sl, "Analysis 4: Driver Win Rate — Who Dominated Their Era?")

add_image_wide(sl, os.path.join(BASE, "pres_winrate.png"), w=9.4, h=3.5)

info_cards_bottom(sl, [
    ("Fangio: 47% win rate (51 starts)",
     "Greatest within-era dominance in history"),
    ("Modern era: lower %, higher totals",
     "20+ race calendars dilute individual win rate"),
    ("Win rate vs. total wins: complementary",
     "Rate = era dominance; total = career longevity"),
    ("Minimum 50 starts filter applied",
     "Removes statistical noise from short careers"),
])


# ════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Conclusions
# ════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
set_bg(sl, DARK)

# Top bar
add_rect(sl, 0, 0, 10, 1.1, RED)
add_text(sl, "Conclusions", 0.4, 0.1, 9.0, 0.9, 32, bold=True, colour=WHITE)

conclusion_cards = [
    ("Geographic expansion mirrors commercialisation",
     "Europe dominated until the 1990s; Middle East & Asia entered from 2000 onward"),
    ("Half of all race entries end in a finish",
     "Reliability improved dramatically from the V8 era; mechanical failures nearly halved"),
    ("Pit stop strategy varies with grid position",
     "Front-starters use ~1.5 stops; mid-field drivers average 2+ as they try to gain places"),
    ("Win rate separates era dominance from longevity",
     "Fangio's 47% remains unmatched; modern champions accumulate wins over a longer calendar"),
]

card_w, card_h = 4.4, 1.6
positions = [(0.5, 1.25), (5.1, 1.25), (0.35, 3.1), (5.1, 3.1)]
for (l, t), (hdr, sub) in zip(positions, conclusion_cards):
    add_rect(sl, l, t, card_w, card_h, CARD)
    add_rect(sl, l, t, card_w, 0.1, RED)    # top accent line
    add_text(sl, hdr, l + 0.15, t + 0.15, card_w - 0.3, 0.5,
             10, bold=True, colour=WHITE, wrap=True)
    add_text(sl, sub, l + 0.15, t + 0.7,  card_w - 0.3, 0.8,
             8.5, bold=False, colour=MUTED, wrap=True)

add_text(sl, "Thank you  ·  Questions?",
         0, 5.15, 10, 0.4, 12, bold=True, colour=RED, align=PP_ALIGN.CENTER)


# ── save ──────────────────────────────────────────────────────────────────────
out = os.path.join(BASE, "F1_Analysis_Presentation.pptx")
prs.save(out)
print(f"Saved → {out}")
